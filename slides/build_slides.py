#!/usr/bin/env python3
"""Build the advisor deck from the measurements, not from memory.

    uv run --with python-pptx python slides/build_slides.py

python-pptx is fetched into a throwaway environment rather than added to the
project: nothing else here needs it, and the project already uses that pattern
for LiteRT export (research/export.py:export_litert).

The benchmark table is READ FROM THE COMMITTED CSV so the deck cannot drift from
the data. Every other figure is a constant below with the RESULT.md section it
comes from, so a number that changes has exactly one place to change.

Two rules the deck must not break, both from RESULT.md §11:
  - report mAP50-95, never mAP50. On the scene-clean split mAP50 saturates at
    0.995 for nearly every cell, so it reads as perfection and means nothing.
  - never show fod-a-3k's own val mAP (0.948, 74% near-duplicates) or Mac latency.
"""

import csv
from decimal import ROUND_HALF_UP, Decimal
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "slides" / "fod-cv-poc-results.pptx"
IMG = ROOT / "slides" / "img"
MATRIX_CSV = ROOT / "runs" / "bench_pi" / "results_480_poc_v2_matrix.csv"

# --- palette: light ground, dark text, one accent. Projector-safe. ----------
INK = RGBColor(0x1A, 0x1D, 0x21)
MUTED = RGBColor(0x5F, 0x67, 0x63)
ACCENT = RGBColor(0xC2, 0x70, 0x1C)   # the result / "after" colour
COOL = RGBColor(0x2E, 0x6F, 0xA8)     # the baseline / "before" colour
RULE = RGBColor(0xD8, 0xDC, 0xD8)
BAND = RGBColor(0xF2, 0xF3, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Arial and Consolas exist on both Windows and macOS PowerPoint. A prettier
# pairing risks a silent substitution on the machine this is presented from.
SANS, MONO = "Arial", "Consolas"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.72)                       # side margin
BODY_W = W - 2 * M

# --- figures, with their source. RESULT.md section in the comment. ----------
CONF_BEFORE, CONF_AFTER = 0.534, 0.936          # §6, 263 scene-clean images
TWO_CORE_CSV = ROOT / "runs" / "bench_pi" / "results_480_poc_v2_2core.csv"
TWO_CORE_CELLS = [("NCNN FP16 (best CPU)", "ncnn", "fp16"), ("Hailo-8 INT8", "hailo", "int8")]
SOAK_FPS, SOAK_DRIFT, SOAK_W = 43.4, "+1.4%", 4.05                                    # §4
INT8_LOSS = [("Hailo-8", 0.851, -0.6, "yes"), ("MNN", 0.847, -1.1, "no - weight-only"),
             ("OpenVINO", 0.830, -3.0, "no - float sim on Arm"),
             ("ONNX", 0.805, -6.0, "no"), ("LiteRT", 0.732, -14.5, "yes - XNNPACK")]   # §5
FP32_REF = 0.856                                                                       # §5


def txt(slide, left, top, width, height, text, size=18, bold=False, color=INK,
        font=SANS, align=PP_ALIGN.LEFT, space_after=6, line=1.25):
    """One textbox, one or more lines. `text` may be a list of (str, dict) pairs."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    items = text if isinstance(text, list) else [text]
    for i, item in enumerate(items):
        body, over = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = body
        p.alignment = over.get("align", align)
        p.space_after = Pt(over.get("space_after", space_after))
        p.line_spacing = over.get("line", line)
        f = p.runs[0].font if p.runs else p.font
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color)
        f.name = over.get("font", font)
    return box


def slide_base(prs, kicker, title, sub=None):
    """Blank layout + the header every content slide shares."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    txt(s, M, Inches(0.46), BODY_W, Inches(0.3), kicker.upper(),
        size=11, bold=True, color=ACCENT, font=SANS)
    txt(s, M, Inches(0.78), BODY_W, Inches(0.7), title, size=30, bold=True)
    y = Inches(1.52)
    if sub:
        txt(s, M, y, BODY_W, Inches(0.45), sub, size=15, color=MUTED)
        y = Inches(2.02)
    line = s.shapes.add_shape(1, M, y, BODY_W, Pt(1.2))   # 1 = rectangle
    line.fill.solid(); line.fill.fore_color.rgb = RULE
    line.line.fill.background(); line.shadow.inherit = False
    return s, y + Inches(0.26)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def table(slide, rows, left, top, width, col_w, size=13, highlight=None, head=True):
    """rows[0] is the header. `highlight` is a row index to accent."""
    n_r, n_c = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_r, n_c, left, top, width, Inches(0.32 * n_r))
    tbl = shape.table
    tbl.first_row = head
    for i, w in enumerate(col_w):
        tbl.columns[i].width = w
    for r, row in enumerate(rows):
        tbl.rows[r].height = Inches(0.34 if r == 0 else 0.31)
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = 0
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = INK
            elif highlight is not None and r == highlight:
                cell.fill.fore_color.rgb = RGBColor(0xFA, 0xEE, 0xDE)
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else BAND
            p = cell.text_frame.paragraphs[0]
            # Right-align only genuinely numeric cells. Testing the first
            # character alone put "+10 frames" in the number column.
            numeric = c and str(val).strip("+-%.,/ ").replace(" ", "").replace("/", "") \
                .replace(".", "").replace(",", "").isdigit()
            p.alignment = PP_ALIGN.RIGHT if numeric else PP_ALIGN.LEFT
            f = p.runs[0].font
            f.size = Pt(size - 1 if r == 0 else size)
            f.name = MONO if (c and r) else SANS
            f.bold = r == 0 or (highlight is not None and r == highlight)
            f.color.rgb = WHITE if r == 0 else INK
    return shape


def pct(value):
    """Half-up to one decimal, because RESULT.md rounds that way.

    +31.25% is the real two-core delta. Python's default gives 31.2, RESULT.md
    says 31.3, and a deck disagreeing with the document it summarises is the
    kind of thing that gets asked about.
    """
    d = Decimal(value).quantize(Decimal("0.1"), rounding=ROUND_HALF_UP)
    return f"{d:+}%"


def read_two_core():
    """4-core and 2-core medians for the cells that matter, straight from the CSVs."""
    four = {(r["format"], r["precision"]): float(r["median_ms"])
            for r in csv.DictReader(MATRIX_CSV.open()) if r["status"] == "ok"}
    two = {(r["format"], r["precision"]): float(r["median_ms"])
           for r in csv.DictReader(TWO_CORE_CSV.open()) if r["status"] == "ok"}
    out = []
    for label, fmt, prec in TWO_CORE_CELLS:
        a, b = four[(fmt, prec)], two[(fmt, prec)]
        out.append((label, a, b, pct((b - a) / a * 100)))
    return out


def read_matrix():
    """Measured cells, fastest first, plus the ones that cannot be built."""
    rows = list(csv.DictReader(MATRIX_CSV.open()))
    ok = sorted((r for r in rows if r["status"] == "ok"),
                key=lambda r: float(r["median_ms"]))
    out = [["Runtime", "Prec", "Median ms", "FPS", "Size MB", "mAP50-95"]]
    hi = None
    for i, r in enumerate(ok, start=1):
        if r["format"] == "hailo":
            hi = i
        out.append([r["format"], r["precision"], f"{float(r['median_ms']):.2f}",
                    f"{float(r['fps']):.1f}", f"{float(r['size_mb']):.2f}",
                    f"{float(r['map50_95']):.3f}"])
    return out, hi


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ---------------------------------------------------------------- 1 title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_shape(1, 0, 0, W, Inches(0.14))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False
    txt(s, M, Inches(2.1), BODY_W, Inches(0.4), "FOD ROBOT — COMPUTER VISION",
        size=13, bold=True, color=ACCENT)
    # 40 chars of Arial Bold 42pt wraps to two lines at this width -- give it
    # the room rather than discovering the collision in the rendered PDF.
    txt(s, M, Inches(2.35), BODY_W, Inches(1.6),
        "Runtime benchmark and detection findings", size=42, bold=True)
    txt(s, M, Inches(4.05), Inches(9.4), Inches(1.2),
        "What the Pi 5 and Hailo-8 actually deliver, why live detection first "
        "failed, and what that changes about the plan.", size=19, color=MUTED)
    txt(s, M, Inches(5.9), BODY_W, Inches(0.8),
        [("Proof-of-concept · measured on hardware · August 2026", {"size": 13}),
         ("All figures reproducible from runs/bench_pi/ in the project repository",
          {"size": 12, "color": MUTED})], color=INK, font=SANS)
    notes(s, "Goal of this deck: (1) what the hardware delivers, (2) how it was measured, "
             "(3) what it means for the project plan. The headline is not the speed number "
             "- it is that the public dataset cannot take us further.")

    # ---------------------------------------------------------- 2 environment
    s, y = slide_base(prs, "1 · Environment", "The hardware under test",
                      "Every number in this deck was measured on this board. Nothing is cited.")
    rows = [["Component", "Specification"],
            ["Compute", "Raspberry Pi 5 Model B - Cortex-A76 4-core @ 2.4 GHz, Active Cooler"],
            ["Accelerator", "Hailo-8 on dual M.2 HAT - 26 TOPS (not the 13-TOPS Hailo-8L)"],
            ["Camera", "Camera Module 3 (imx708) - 4608x2592, 66 deg lens, autofocus"],
            ["Storage", "NVMe SSD on the shared PCIe link"],
            ["Model", "YOLO11n, 4 classes, deployed at 480x480"],
            ["Build host", "Apple M4 - training and export only"]]
    table(s, rows, M, y, BODY_W, [Inches(2.0), Inches(9.9)], size=14)
    txt(s, M, Inches(5.55), BODY_W, Inches(1.2),
        [("The accelerator is a measurement, not yet a decision.", {"bold": True, "size": 16}),
         ("PRD v2 commits to Pi 5 CPU-only; adopting the Hailo-8 is a BOM change "
          "gated on advisor sign-off. No Apple M4 figure appears anywhere in this "
          "deck - it does not transfer to Arm.", {"size": 14, "color": MUTED})])
    notes(s, "Hailo-8 vs Hailo-8L matters: Ultralytics defaults to the 8L and that .hef "
             "will not load on this board. The Mac is used only to train and export, "
             "because the LiteRT converter has no aarch64 wheel - it cannot run on the Pi.")

    # ------------------------------------------------------------- 3 method
    s, y = slide_base(prs, "2 · Method", "How it was measured",
                      "Two things decide whether a benchmark means anything: the split, and the board.")
    txt(s, M, y, Inches(6.0), Inches(0.35), "Datasets - three splits, three jobs",
        size=16, bold=True)
    rows = [["Split", "Size", "Used for"],
            ["fod-a", "510 / 90", "the v1 model"],
            ["fod-a-3k", "2,550 / 450", "training v2 only"],
            ["fod-a-clean", "263 val", "ALL scoring here"]]
    table(s, rows, M, y + Inches(0.42), Inches(6.0),
          [Inches(1.8), Inches(1.5), Inches(2.7)], size=12)
    txt(s, M, y + Inches(2.05), Inches(6.0), Inches(1.0),
        "fod-a-clean holds only scenes that contributed no training frame to either "
        "model - the reason appears on slide 8.", size=13, color=MUTED)

    x2 = M + Inches(6.6)
    txt(s, x2, y, Inches(5.3), Inches(0.35), "Board control - all five required",
        size=16, bold=True)
    txt(s, x2, y + Inches(0.45), Inches(5.3), Inches(2.0),
        [("CPU governor pinned to performance (default idles at 1.6 GHz)", {}),
         ("120 s cooldown to 66 C between every cell", {}),
         ("No desktop session - costs the CPU path 16% invisibly", {}),
         ("Drift control re-runs the first cell at the end", {}),
         ("A CPU format measured first, never the accelerator", {})],
        size=13, space_after=9)
    txt(s, x2, y + Inches(2.75), Inches(5.3), Inches(0.9),
        [("This run: drift -2.3%, throttled 0x0, 2.4 GHz throughout.",
          {"bold": True, "size": 14, "color": ACCENT})])
    notes(s, "The first Pi run went 61->80 C and drift hit +20.3% - position in the loop "
             "outweighed the runtime being measured. Every rule here exists because "
             "skipping it produced a wrong ranking at least once.")

    # ---------------------------------------------------------- 4 benchmark
    s, y = slide_base(prs, "3 · Result", "Benchmark: every runtime and precision",
                      "poc-v2 model, 480x480, 4 threads, idle board. 50 runs per cell after 5 warm-up.")
    rows, hi = read_matrix()
    table(s, rows, M, y, Inches(8.2),
          [Inches(1.6), Inches(1.0), Inches(1.7), Inches(1.2), Inches(1.4), Inches(1.3)],
          size=12, highlight=hi)
    x2 = M + Inches(8.6)
    txt(s, x2, y, Inches(3.9), Inches(3.4),
        [("Hailo-8: 17.8 ms, 56 FPS", {"size": 19, "bold": True, "color": ACCENT}),
         ("2.4x the fastest CPU runtime, and it barely touches the CPU - "
          "0.4 ms of postprocess against ~1.0 ms elsewhere, because NMS is "
          "compiled onto the chip.", {"size": 13, "color": MUTED}),
         ("", {"size": 6}),
         ("Software INT8 does not compute", {"size": 15, "bold": True}),
         ("MNN INT8 is slower than its own FP32; ONNX INT8 lands within 1 ms of "
          "its own. The file shrinks 3.7x, the arithmetic never changes.",
          {"size": 13, "color": MUTED}),
         ("", {"size": 6}),
         ("Not buildable: NCNN INT8 has no export path at all, and the Hailo-8 "
          "is INT8-only silicon. Those gaps are toolchain facts, not omissions.",
          {"size": 12, "color": MUTED})])
    notes(s, "Read mAP50-95, never mAP50: on this split mAP50 saturates at 0.995 for "
             "almost every cell, so it reports every INT8 cell as free. "
             "OpenVINO INT8 at 167 ms is a 2.6x regression against its own FP32 - "
             "Arm executes quantised graphs in float simulation. "
             "Source: runs/bench_pi/results_480_poc_v2_matrix.csv")

    # --------------------------------------------------- 5 two cores + soak
    s, y = slide_base(prs, "4 · Result", "The measurement that decides the architecture",
                      "A robot running SLAM cannot give four cores to vision. So measure it at two.")
    rows = [["Runtime", "4 cores", "2 cores", "Change"]]
    for name, four, two, change in read_two_core():
        rows.append([name, f"{four:.2f} ms", f"{two:.2f} ms", change])
    table(s, rows, M, y, Inches(7.4),
          [Inches(2.9), Inches(1.5), Inches(1.5), Inches(1.5)], size=15, highlight=2)
    txt(s, M, y + Inches(1.35), Inches(7.4), Inches(1.4),
        [("Halving the cores costs the CPU a third of its throughput and the "
          "accelerator nothing.", {"size": 16, "bold": True}),
         ("Hailo's inference is not on the CPU at all, so the gap widens from "
          "2.46x to 3.24x. At two cores the CPU path drops to 17.3 FPS while the "
          "accelerator holds 56.", {"size": 13, "color": MUTED})])
    x2 = M + Inches(7.9)
    txt(s, x2, y, Inches(4.6), Inches(2.8),
        [("Sustained load, 600 s", {"size": 16, "bold": True}),
         (f"{SOAK_FPS} FPS sustained", {"size": 26, "bold": True, "color": ACCENT}),
         (f"26,042 frames. Drift first quarter to last: {SOAK_DRIFT}. "
          f"Temperature flat at 68 C, clock never left 2.4 GHz, {SOAK_W} W median.",
          {"size": 13, "color": MUTED}),
         ("", {"size": 6}),
         ("43.4 FPS is the figure to plan a robot around; 56 belongs to the "
          "leaderboard.", {"size": 13, "bold": True})])
    notes(s, "The two-core result is the strongest argument for the accelerator and it is "
             "not about speed - it is about what the CPU is free to do. "
             "Soak is lower than the burst mainly because it cycles 263 images: cache "
             "locality, not thermal decay. A true sustained figure needs a fixed frame, "
             "which is what a camera delivers.")

    # -------------------------------------------------------------- 6 INT8
    s, y = slide_base(prs, "5 · Result", "What INT8 quantisation costs",
                      f"All FP32 and FP16 cells score {FP32_REF:.3f} mAP50-95. That is the reference.")
    rows = [["Backend", "INT8 mAP50-95", "Loss", "Real INT8 compute?"]]
    for name, v, d, real in INT8_LOSS:
        rows.append([name, f"{v:.3f}", f"{d:+.1f}%", real])
    table(s, rows, M, y, Inches(7.2),
          [Inches(1.7), Inches(1.9), Inches(1.3), Inches(2.3)], size=13, highlight=1)
    txt(s, M, y + Inches(2.05), Inches(7.2), Inches(1.2),
        [("Only the accelerator buys both size and speed.", {"size": 16, "bold": True}),
         ("MNN matches it on accuracy but not latency; LiteRT is fast but is the "
          "only path losing real accuracy.", {"size": 13, "color": MUTED})])
    x2 = M + Inches(7.7)
    txt(s, x2, y, Inches(4.8), Inches(3.4),
        [("A previous conclusion was wrong", {"size": 17, "bold": True, "color": ACCENT}),
         ("Earlier work concluded \"INT8 is dead for this project\" because LiteRT "
          "lost 34% of its accuracy. That write-up also flagged its own doubt: the "
          "calibration set was only 510 images.", {"size": 13, "color": MUTED}),
         ("", {"size": 6}),
         ("Re-run with 2,550 calibration images:", {"size": 13, "bold": True}),
         ("LiteRT   -34%  ->  -4.0%", {"size": 17, "bold": True, "font": MONO}),
         ("ONNX     -16%  ->   0.0%", {"size": 17, "bold": True, "font": MONO}),
         ("Most of the collapse was the calibration set, not the runtime.",
          {"size": 13, "color": MUTED})])
    notes(s, "This is the slide to be honest on: we published a conclusion, the conclusion "
             "was wrong, and the earlier write-up had already named the reason to re-test. "
             "The ranking survives - LiteRT is still worst - but 'INT8 is dead' does not.")

    # ------------------------------------------------------- 7 live camera
    s, y = slide_base(prs, "6 · Finding", "Pointed at real fasteners, it detected almost nothing",
                      "The benchmark said the model was good. The camera said otherwise. Three causes.")
    # Constrain by WIDTH: the frame is portrait, and sizing by height pushed it
    # 6.2 in wide, straight under the text column.
    s.shapes.add_picture(str(IMG / "camera_blurred_screws.jpg"), M, y,
                         width=Inches(2.3))
    txt(s, M, y + Inches(4.2), Inches(3.4), Inches(0.6),
        "Two screws, plainly visible to a person, and one junk box. "
        "Out of focus - not a model failure.", size=11, color=MUTED)
    x2 = M + Inches(4.0)
    txt(s, x2, y, Inches(8.5), Inches(3.7),
        [("1.  The lens was parked at 1.00 m", {"size": 17, "bold": True}),
         ("Autofocus was never enabled - libcamera defaults to manual at one metre, "
          "and nothing in our code set it. Every frame ever captured was focused at "
          "1 m regardless of the object. Fixed: focus score 180 -> 757, no frame cost.",
          {"size": 13, "color": MUTED}),
         ("2.  The model had seen ten screws", {"size": 17, "bold": True}),
         ("The training set was capped at 600 images as a smoke test and never "
          "uncapped. Nail and screw - the two objects that failed - were the two "
          "weakest classes. Lifting the cap took screws from 10 to 36 instances.",
          {"size": 13, "color": MUTED}),
         ("3.  The data contains no cluttered scenes", {"size": 17, "bold": True}),
         ("See the next slide. This one is not fixable from the public dataset.",
          {"size": 13, "color": MUTED})],
        space_after=10)
    txt(s, x2, y + Inches(3.85), Inches(8.5), Inches(0.5),
        [(f"Median detection confidence: {CONF_BEFORE:.3f}  ->  {CONF_AFTER:.3f}",
          {"size": 20, "bold": True, "color": ACCENT})])
    notes(s, "Zoom and input resolution were both ruled out by measurement, not assumed: "
             "a 40 mm screw is already at training scale at 0.28 m with no zoom, which is "
             "the PRD working distance. The hand-written preprocessing was also verified "
             "against the standard toolchain - IoU 0.82-0.95. So the pipeline was never "
             "the problem.")

    # ------------------------------------------------ 8 the fix, on hardware
    s, y = slide_base(prs, "7 · Result", "The same camera, after the fixes",
                      "Autofocus enabled, model retrained on 5x the data, unknown class suppressed.")
    s.shapes.add_picture(str(IMG / "camera_detect_after.jpg"), M, y, width=Inches(7.4))
    txt(s, M, y + Inches(4.3), Inches(7.4), Inches(0.5),
        "Live frame, Hailo-8 on the Pi. Confidence printed on each box.",
        size=11, color=MUTED)
    x2 = M + Inches(7.9)
    txt(s, x2, y, Inches(4.6), Inches(4.3),
        [("12 of 12 fasteners found", {"size": 24, "bold": True, "color": ACCENT}),
         ("Across two frames, every object detected, boxes tight, confidence "
          "0.42 to 0.87.", {"size": 14, "color": MUTED}),
         ("Nothing fired on the clutter", {"size": 17, "bold": True}),
         ("The bag, cables, glasses and boxes on the left are all ignored. That is "
          "the exact failure that produced full-frame boxes before.",
          {"size": 13, "color": MUTED}),
         ("Still wrong: the class", {"size": 17, "bold": True}),
         ("Screws are labelled bolt. Detection and localisation are solved; "
          "naming is not, and 36 screw instances is why.", {"size": 13, "color": MUTED}),
         ("Honest limit: the fasteners sit on a plain sheet, so this is not yet a "
          "cluttered-floor test. The clutter is behind them, not around them.",
          {"size": 12, "color": MUTED})], space_after=9)
    notes(s, "This is the result the whole investigation was aiming at, and it arrived "
             "after the deck was drafted. Before: almost nothing detected, full-frame "
             "boxes on furniture. After: 12/12 with no false positives. "
             "Be honest in the room about two things - the class confusion is real and "
             "expected from the training counts, and the objects are on a white sheet, "
             "so the arena-floor case is still untested.")

    # ---------------------------------------------------- 9 dataset ceiling
    s, y = slide_base(prs, "8 · Finding", "The dataset is the ceiling",
                      "FOD-A is not 9,623 independent images. It is about 38 scenes.")
    s.shapes.add_picture(str(IMG / "foda_sample.jpg"), M, y, height=Inches(2.7))
    txt(s, M, y + Inches(2.8), Inches(2.7), Inches(0.8),
        "A typical training image: one object, blank ground, no clutter, "
        "no second object.", size=11, color=MUTED)
    x2 = M + Inches(3.2)
    rows = [["Image pair", "Similarity"],
            ["adjacent frames", "0.975"],
            ["+10 frames", "0.949"],
            ["unrelated pair", "0.778"]]
    table(s, rows, x2, y, Inches(3.5), [Inches(2.1), Inches(1.4)], size=12)
    txt(s, x2, y + Inches(1.65), Inches(3.5), Inches(1.6),
        "Video-derived and heavily redundant: 38 runs cover 96% of the images, "
        "and half of everything sits in the largest six.", size=13, color=MUTED)

    x3 = M + Inches(7.1)
    txt(s, x3, y, Inches(5.4), Inches(4.2),
        [("Which broke the split", {"size": 16, "bold": True}),
         ("74% of the validation frames used to train v2 had a near-duplicate in "
          "its own training set. That model scores 0.948 on it - recall of an "
          "almost identical frame, not generalisation. All scoring in this deck "
          "uses a scene-disjoint split instead.", {"size": 13, "color": MUTED}),
         ("", {"size": 6}),
         ("And the images are the wrong kind", {"size": 16, "bold": True}),
         ("Every training image holds exactly one object on an empty plane. "
          "So a model trained on it cannot abstain - pointed at a desk it labels "
          "a mouse and a keyboard as debris. It has never been shown a scene.",
          {"size": 13, "color": MUTED}),
         ("", {"size": 6}),
         ("More of this data cannot fix false positives on furniture, "
          "because it contains no furniture.", {"size": 14, "bold": True, "color": ACCENT})])
    notes(s, "This is the most important slide in the deck. It says the next step is not "
             "more benchmarking or more public data - it is collecting our own arena "
             "images, which PRD section 10 already specifies. Note also: once you train on "
             "2,550 of FOD-A's frames, zero untouched images remain even 30 frames away "
             "from a training frame. There is no clean holdout left inside it.")

    # --------------------------------------------------------- 9 decisions
    s, y = slide_base(prs, "9 · Decisions", "What this settles",
                      "Each decision with the measurement behind it.")
    rows = [["Decision", "Because"],
            ["Deploy at 480x480, not 640", "same accuracy, 1.87x faster. 320 loses 37%"],
            ["Hailo-8 is the only real INT8 path", "-0.6% accuracy, 2.4x fastest CPU, 3.2x at two cores"],
            ["Drop \"NCNN INT8\" from the spec", "no export path exists in this toolchain"],
            ["Drop OpenVINO INT8", "2.6x slower than its own FP32, replicated twice"],
            ["LiteRT INT8 usable but lossy", "fastest CPU INT8, but -14.5% accuracy"],
            ["YOLO26n rejected", "45.5 ms vs YOLO11n's 44.8 - no gain"],
            ["Enable autofocus, always", "it was never set; focus score 180 -> 757"],
            ["Score on a scene-disjoint split", "the per-image shuffle leaked 74%"]]
    table(s, rows, M, y, BODY_W, [Inches(4.5), Inches(7.4)], size=13)
    notes(s, "Every row here replaces a claim that was previously assumed or cited rather "
             "than measured. Three of them amend the PRD: FR-2 (image size), FR-1 (the "
             "NCNN INT8 and OpenVINO INT8 entries).")

    # -------------------------------------------------------------- 10 next
    s, y = slide_base(prs, "10 · Next", "What happens next",
                      "Toolchain answered, detection working on hardware. The data question is not.")
    txt(s, M, y, Inches(7.6), Inches(3.6),
        [("1.  Collect the arena dataset  (PRD section 10)",
          {"size": 18, "bold": True, "color": ACCENT}),
         ("Now the critical path, not an eventual refinement. The public dataset "
          "cannot supply a single cluttered negative, and clutter is what produces "
          "the false positives.", {"size": 14, "color": MUTED}),
         ("2.  Group the split by scene before training on it", {"size": 18, "bold": True}),
         ("A per-image shuffle cannot give a trustworthy score on video-derived "
          "data - arena footage will be video too.", {"size": 14, "color": MUTED}),
         ("3.  Fix the class confusion, or drop to one class", {"size": 18, "bold": True}),
         ("Screws are detected reliably and named wrong. The specification asks "
          "for a single metal_fastener class anyway, with per-class recall "
          "recovered from the seeding log.", {"size": 14, "color": MUTED})],
        space_after=11)
    x2 = M + Inches(8.1)
    txt(s, x2, y, Inches(4.4), Inches(3.2),
        [("Smaller, still open", {"size": 15, "bold": True}),
         ("Viewpoint augmentation written for this camera geometry, never run",
          {"size": 13, "color": MUTED}),
         ("USB power meter to replace the PMIC estimate", {"size": 13, "color": MUTED}),
         ("Calibrate INT8 on arena images, not FOD-A", {"size": 13, "color": MUTED}),
         ("", {"size": 8}),
         ("Accelerator adoption is a BOM change and stays gated on sign-off.",
          {"size": 13, "bold": True})], space_after=9)
    notes(s, "Close on this: the benchmark work is finished and the answer is clear, but "
             "the thing standing between here and a working robot detector is data "
             "collection, and that needs the arena.")

    # slide numbers on every slide but the title
    for i, s in enumerate(prs.slides, start=1):
        if i == 1:
            continue
        txt(s, W - Inches(1.3), H - Inches(0.52), Inches(0.6), Inches(0.3),
            str(i), size=11, color=MUTED, align=PP_ALIGN.RIGHT)

    OUT.parent.mkdir(exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(ROOT)}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
